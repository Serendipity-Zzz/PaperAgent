import csv
import hashlib
import json
import zipfile
from pathlib import Path

import nbformat
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from paperagent.preview.renderers import (
    ArchiveRenderer,
    StructuredDocumentRenderer,
    TableRenderer,
)


def sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def test_structured_office_email_bibliography_notebook_and_json(tmp_path: Path) -> None:
    docx = tmp_path / "论文.docx"
    document = Document()
    document.add_heading("标题", level=1)
    document.add_table(rows=1, cols=2).rows[0].cells[0].text = "复杂表格"
    document.save(docx)

    pptx = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "实验"
    slide.notes_slide.notes_text_frame.text = "演讲者备注"
    presentation.save(pptx)

    email = tmp_path / "mail.eml"
    email.write_text(
        "Subject: Review\nMessage-ID: <review@test>\n"
        "Date: Wed, 1 Jan 2025 12:00:00 +0800\n\nRevise it",
        encoding="utf-8",
    )
    bib = tmp_path / "refs.bib"
    bib.write_text("@article{x,title={Local Agent},year={2026}}", encoding="utf-8")
    notebook = tmp_path / "unsafe.ipynb"
    nbformat.write(
        nbformat.v4.new_notebook(
            cells=[
                nbformat.v4.new_code_cell(
                    "raise RuntimeError('must not run')",
                    outputs=[
                        nbformat.v4.new_output(
                            "display_data", data={"text/html": "<script>x()</script>"}
                        )
                    ],
                )
            ]
        ),
        notebook,
    )
    structured_json = tmp_path / "chat.json"
    structured_json.write_text(
        json.dumps([{"id": "m1", "participant": "user", "content": "Decision"}]),
        encoding="utf-8",
    )

    renderer = StructuredDocumentRenderer()
    for path in (docx, pptx, email, bib, notebook, structured_json):
        result = renderer.render(path, file_id=path.name, source_hash=sha(path))
        assert result.parts, path.name
        assert all(part.anchor is not None for part in result.parts)
    notebook_result = renderer.render(notebook, file_id="nb", source_hash=sha(notebook))
    assert "must not run" in str(notebook_result.parts[0].payload)
    assert "<script>" not in str(notebook_result.parts[0].payload)


def test_tables_are_paged_by_service_contract_and_keep_sheet_cell_anchors(tmp_path: Path) -> None:
    csv_path = tmp_path / "十万行.csv"
    with csv_path.open("w", encoding="utf-8", newline="") as stream:
        writer = csv.writer(stream)
        writer.writerows([index, f"value-{index}"] for index in range(100_000))
    result = TableRenderer().render(csv_path, file_id="csv", source_hash=sha(csv_path))
    assert result.payload["row_count"] == 100_000
    assert result.payload["virtualized"] is True
    assert result.parts[-1].anchor.cell_range == "A100000"

    xlsx = tmp_path / "公式与空表.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    sheet.merge_cells("A1:B1")
    sheet["A1"] = "标题"
    sheet["A2"] = "=1/0"
    workbook.create_sheet("空表")
    workbook.save(xlsx)
    excel = TableRenderer().render(xlsx, file_id="xlsx", source_hash=sha(xlsx))
    assert excel.payload["sheets"] == ["数据", "空表"]
    assert excel.parts[0].anchor.sheet == "数据"


def test_archive_tree_never_extracts_and_rejects_traversal(tmp_path: Path) -> None:
    safe = tmp_path / "safe.zip"
    with zipfile.ZipFile(safe, "w") as archive:
        archive.writestr("folder/readme.txt", "safe")
    result = ArchiveRenderer().render(safe, file_id="zip", source_hash=sha(safe))
    assert result.payload["extraction_disabled"] is True
    assert not (tmp_path / "folder").exists()

    bad = tmp_path / "bad.zip"
    with zipfile.ZipFile(bad, "w") as archive:
        archive.writestr("../escape.txt", "bad")
    try:
        ArchiveRenderer().render(bad, file_id="bad", source_hash=sha(bad))
    except ValueError as error:
        assert "unsafe path" in str(error)
    else:
        raise AssertionError("path traversal archive was accepted")
