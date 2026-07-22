import json
import mailbox
import zipfile
from email.message import EmailMessage
from pathlib import Path

import nbformat
import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

from paperagent.ingestion.parsers import default_registry


def test_registry_text_duplicate_cancel_and_empty(tmp_path: Path) -> None:
    registry = default_registry()
    path = tmp_path / "需求.md"
    path.write_text("# 目标\n必须离线工作", encoding="utf-8")
    first = registry.import_file(path)
    duplicate = registry.import_file(path)
    assert len(first.source.chunks) == 2
    assert first.source.chunks[0].kind == "heading"
    assert first.source.chunks[0].text == "目标"
    assert duplicate.duplicate_of == first.source.id
    other = tmp_path / "cancel.txt"
    other.write_text("content", encoding="utf-8")
    assert registry.import_file(other, cancelled=lambda: True).cancelled
    empty = tmp_path / "empty.txt"
    empty.touch()
    with pytest.raises(ValueError, match="empty"):
        registry.import_file(empty)


def test_safe_html_docx_pdf_and_tables(tmp_path: Path) -> None:
    registry = default_registry()
    html = tmp_path / "unsafe.html"
    html.write_text("<script>alert(1)</script><h1>Safe title</h1>", encoding="utf-8")
    html_report = registry.import_file(html)
    assert html_report.source.chunks[0].text == "Safe title"
    assert html_report.warnings

    docx = tmp_path / "sample.docx"
    word = Document()
    word.add_heading("第一章 1. Heading", level=1)
    word.add_paragraph("Paragraph")
    word.save(docx)
    docx_report = registry.import_file(docx)
    assert len(docx_report.source.chunks) == 2
    assert docx_report.source.chunks[0].text == "Heading"

    xlsx = tmp_path / "table.xlsx"
    workbook = Workbook()
    workbook.active.append(["name", "value"])
    workbook.save(xlsx)
    assert registry.import_file(xlsx).source.chunks[0].kind == "table_row"

    pdf = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with pdf.open("wb") as stream:
        writer.write(stream)
    assert "OCR" in registry.import_file(pdf).warnings[0]


def test_bibliography_email_code_notebook_and_archive_security(tmp_path: Path) -> None:
    registry = default_registry()
    bib = tmp_path / "refs.bib"
    bib.write_text("@article{x, title={Paper}, doi={10.1/test}}", encoding="utf-8")
    assert registry.import_file(bib).source.chunks[0].metadata["doi"] == "10.1/test"

    eml = tmp_path / "mail.eml"
    eml.write_text("Subject: Decision\nMessage-ID: <one@test>\n\nUse SQLite", encoding="utf-8")
    assert registry.import_file(eml).source.chunks[0].citation_policy == "internal_only"

    code = tmp_path / "module.py"
    code.write_text("def build():\n    return 1", encoding="utf-8")
    assert registry.import_file(code).source.chunks[0].kind == "symbol"

    archive = tmp_path / "bad.zip"
    with zipfile.ZipFile(archive, "w") as output:
        output.writestr("../escape.txt", "bad")
    with pytest.raises(ValueError, match="traversal"):
        registry.import_file(archive)


def test_communications_notebook_images_endnote_and_directory_filters(tmp_path: Path) -> None:
    registry = default_registry()
    chat = tmp_path / "chat.json"
    chat.write_text(
        json.dumps(
            [
                {
                    "id": "m1",
                    "participant": "user",
                    "timestamp": "2026-01-01",
                    "content": "Decision",
                },
                {"id": "m2", "participant": "agent", "content": "Action item"},
            ]
        ),
        encoding="utf-8",
    )
    chat_report = registry.import_file(chat)
    assert len(chat_report.source.chunks) == 2
    assert all(chunk.citation_policy == "process_only" for chunk in chat_report.source.chunks)

    notebook_path = tmp_path / "analysis.ipynb"
    notebook = nbformat.v4.new_notebook(
        cells=[nbformat.v4.new_code_cell("print('must not execute')")]
    )
    nbformat.write(notebook, notebook_path)
    notebook_report = registry.import_file(notebook_path)
    assert notebook_report.source.chunks[0].text == "print('must not execute')"

    image_path = tmp_path / "design.png"
    Image.new("RGB", (32, 24), "white").save(image_path)
    assert registry.import_file(image_path).source.metadata["width"] == 32
    svg = tmp_path / "unsafe.svg"
    svg.write_text("<svg><script>alert(1)</script></svg>", encoding="utf-8")
    with pytest.raises(ValueError, match="active"):
        registry.import_file(svg)

    endnote = tmp_path / "endnote.xml"
    endnote.write_text(
        "<xml><records><record><titles><title>Paper</title></titles>"
        "<electronic-resource-num>10.1/example</electronic-resource-num></record></records></xml>",
        encoding="utf-8",
    )
    assert registry.import_file(endnote).source.chunks[0].metadata["doi"] == "10.1/example"

    mbox_path = tmp_path / "mail.mbox"
    box = mailbox.mbox(mbox_path)
    message = EmailMessage()
    message["Message-ID"] = "<one@example.test>"
    message["Subject"] = "Meeting"
    message.set_content("Decision and action")
    box.add(message)
    box.flush()
    box.close()
    assert registry.import_file(mbox_path).source.chunks[0].kind == "email"

    presentation_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Research design"
    presentation.save(presentation_path)
    assert registry.import_file(presentation_path).source.chunks[0].kind == "slide"

    repository = tmp_path / "repo"
    repository.mkdir()
    (repository / "main.py").write_text("def main(): pass", encoding="utf-8")
    ignored = repository / "node_modules"
    ignored.mkdir()
    (ignored / "secret.py").write_text("password='hidden'", encoding="utf-8")
    reports = registry.import_directory(repository)
    assert len(reports) == 1 and reports[0].source.name == "main.py"
    assert reports[0].source.metadata["relative_path"] == "main.py"
