from __future__ import annotations

import csv
import email
import json
import mailbox
import re
import zipfile
from email import policy
from pathlib import Path
from typing import cast

import bibtexparser
import extract_msg
import nbformat
import rispy
from bs4 import BeautifulSoup
from bs4.element import Tag
from charset_normalizer import from_bytes
from defusedxml import ElementTree
from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from paperagent.ingestion.registry import IngestionRegistry, Parser
from paperagent.ingestion.schemas import (
    Chunk,
    CitationPolicy,
    ImportReport,
    Locator,
    SourceDocument,
)
from paperagent.schemas.numbering import NumberingNormalizer

_NUMBERING_NORMALIZER = NumberingNormalizer()


def source(path: Path, digest: str, parser: str, media_type: str) -> SourceDocument:
    return SourceDocument(name=path.name, sha256=digest, parser=parser, media_type=media_type)


class TextParser:
    name = "text"
    extensions = (".txt", ".md", ".markdown")

    def parse(self, path: Path, digest: str) -> ImportReport:
        raw = path.read_bytes()
        try:
            decoded = raw.decode("utf-8")
        except UnicodeDecodeError:
            match = from_bytes(raw).best()
            if match is None:
                raise ValueError("Unable to detect text encoding") from None
            decoded = str(match)
        document = source(
            path, digest, self.name, "text/markdown" if path.suffix != ".txt" else "text/plain"
        )
        for line_number, line in enumerate(decoded.splitlines(), start=1):
            if line.strip():
                kind = "text"
                chunk_text = line
                if path.suffix in {".md", ".markdown"}:
                    heading_match = re.match(r"^#{1,6}\s+(.+)$", line)
                    if heading_match:
                        kind = "heading"
                        chunk_text = _NUMBERING_NORMALIZER.normalize(
                            heading_match.group(1), node_kind="heading"
                        ).semantic
                    elif re.search(r"\$\$?.+?\$\$?", line):
                        kind = "equation"
                    elif re.search(r"!\[[^]]*]\([^)]+\)", line):
                        kind = "image_reference"
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=chunk_text,
                        kind=kind,
                        locator=Locator(line_start=line_number, line_end=line_number),
                    )
                )
        return ImportReport(source=document)


class HtmlParser:
    name = "html"
    extensions = (".html", ".htm")

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "text/html")
        soup = BeautifulSoup(path.read_bytes(), "html.parser")
        dangerous = soup.find_all(["script", "iframe", "object", "embed"])
        for node in dangerous:
            node.decompose()
        for index, node in enumerate(soup.find_all(["h1", "h2", "h3", "p", "li", "td"])):
            text = node.get_text(" ", strip=True)
            if text:
                document.chunks.append(
                    Chunk(
                        source_id=document.id, text=text, locator=Locator(json_path=f"dom[{index}]")
                    )
                )
        for index, node in enumerate(soup.find_all(["img", "math"])):
            if not isinstance(node, Tag):
                continue
            if node.name == "img":
                media_text = str(node.get("alt") or node.get("src") or f"image-{index}")
                kind = "image_reference"
            else:
                media_text = node.get_text(" ", strip=True)
                kind = "equation"
            document.chunks.append(
                Chunk(
                    source_id=document.id,
                    text=media_text,
                    kind=kind,
                    locator=Locator(json_path=f"media[{index}]"),
                )
            )
        warnings = [f"Removed {len(dangerous)} active HTML element(s)"] if dangerous else []
        return ImportReport(source=document, warnings=warnings)


class PdfParser:
    name = "pdf"
    extensions = (".pdf",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/pdf")
        warnings: list[str] = []
        reader = PdfReader(path)
        for page_number, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                document.chunks.append(
                    Chunk(source_id=document.id, text=text, locator=Locator(page=page_number))
                )
                heading_candidates: list[dict[str, object]] = []
                for line in (item.strip() for item in text.splitlines() if item.strip()):
                    result = _NUMBERING_NORMALIZER.dry_run(line, node_kind="heading")
                    if result.changed:
                        heading_candidates.append(
                            {
                                "original": result.original,
                                "semantic": result.semantic,
                                "prefixes": [
                                    prefix.model_dump(mode="json")
                                    for prefix in result.prefixes
                                ],
                            }
                        )
                if heading_candidates:
                    candidates = cast(
                        dict[str, list[dict[str, object]]],
                        document.metadata.setdefault("normalized_heading_candidates", {}),
                    )
                    candidates[str(page_number)] = heading_candidates
            else:
                warnings.append(f"Page {page_number} has no text layer; OCR may be required")
            try:
                image_count = len(page.images)
            except (AttributeError, TypeError):
                image_count = 0
            if image_count:
                page_images = cast(dict[str, int], document.metadata.setdefault("page_images", {}))
                page_images[str(page_number)] = image_count
        return ImportReport(source=document, warnings=warnings)


class DocxParser:
    name = "docx"
    extensions = (".docx",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(
            path,
            digest,
            self.name,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        word = Document(str(path))
        for index, paragraph in enumerate(word.paragraphs):
            if paragraph.text.strip():
                style_name = paragraph.style.name if paragraph.style else ""
                kind = "heading" if style_name.startswith("Heading") else "paragraph"
                paragraph_text = paragraph.text
                if kind == "heading":
                    paragraph_text = _NUMBERING_NORMALIZER.normalize(
                        paragraph_text, node_kind="heading"
                    ).semantic
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=paragraph_text,
                        kind=kind,
                        locator=Locator(paragraph=index),
                        metadata={"style": style_name},
                    )
                )
                equation_nodes = paragraph._p.xpath(".//m:oMath")
                for equation_index, equation in enumerate(equation_nodes):
                    document.chunks.append(
                        Chunk(
                            source_id=document.id,
                            text="".join(equation.itertext()),
                            kind="equation",
                            locator=Locator(
                                json_path=f"paragraphs[{index}].equations[{equation_index}]"
                            ),
                        )
                    )
                image_nodes = paragraph._p.xpath(".//a:blip")
                for image_index, image in enumerate(image_nodes):
                    relationship_id = image.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    document.chunks.append(
                        Chunk(
                            source_id=document.id,
                            text=relationship_id or f"image-{image_index}",
                            kind="image_reference",
                            locator=Locator(json_path=f"paragraphs[{index}].images[{image_index}]"),
                        )
                    )
        for table_index, table in enumerate(word.tables):
            for row_index, row in enumerate(table.rows):
                text = "\t".join(cell.text for cell in row.cells)
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=text,
                        kind="table_row",
                        locator=Locator(json_path=f"tables[{table_index}].rows[{row_index}]"),
                    )
                )
        return ImportReport(source=document)


class PptxParser:
    name = "pptx"
    extensions = (".pptx",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(
            path,
            digest,
            self.name,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        presentation = Presentation(str(path))
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if slide.has_notes_slide:
                notes = [
                    shape.text
                    for shape in slide.notes_slide.notes_text_frame.paragraphs
                    if getattr(shape, "text", "").strip()
                ]
                texts.extend(notes)
            if texts:
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text="\n".join(texts),
                        kind="slide",
                        locator=Locator(json_path=f"slides[{slide_number - 1}]"),
                        metadata={"slide_number": slide_number},
                    )
                )
        return ImportReport(source=document)


class CsvParser:
    name = "csv"
    extensions = (".csv", ".tsv")

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "text/csv")
        match = from_bytes(path.read_bytes()).best()
        if match is None:
            raise ValueError("Unable to decode table")
        delimiter = "\t" if path.suffix == ".tsv" else ","
        for row_index, row in enumerate(
            csv.reader(str(match).splitlines(), delimiter=delimiter), start=1
        ):
            document.chunks.append(
                Chunk(
                    source_id=document.id,
                    text="\t".join(row),
                    kind="table_row",
                    locator=Locator(sheet="Sheet1", cell_range=f"A{row_index}"),
                )
            )
        return ImportReport(source=document)


class XlsxParser:
    name = "xlsx"
    extensions = (".xlsx", ".xlsm")

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(
            path,
            digest,
            self.name,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        workbook = load_workbook(path, read_only=True, data_only=False, keep_vba=False)
        for sheet in workbook.worksheets:
            merged_cells = getattr(sheet, "merged_cells", None)
            merged_ranges = [str(item) for item in merged_cells.ranges] if merged_cells else []
            for row in sheet.iter_rows():
                values = ["" if cell.value is None else str(cell.value) for cell in row]
                if any(values):
                    document.chunks.append(
                        Chunk(
                            source_id=document.id,
                            text="\t".join(values),
                            kind="table_row",
                            locator=Locator(sheet=sheet.title, cell_range=f"A{row[0].row}"),
                            metadata={
                                "formula_cells": [
                                    cell.coordinate
                                    for cell in row
                                    if isinstance(cell.value, str) and cell.value.startswith("=")
                                ],
                                "merged_ranges": merged_ranges,
                            },
                        )
                    )
        workbook.close()
        return ImportReport(source=document)


class BibliographyParser:
    name = "bibliography"
    extensions = (".bib", ".ris")

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/x-research-info-systems")
        text = path.read_text(encoding="utf-8", errors="replace")
        entries = rispy.loads(text) if path.suffix == ".ris" else bibtexparser.loads(text).entries
        for index, entry in enumerate(entries):
            title = entry.get("title") or entry.get("primary_title") or "Untitled"
            doi = entry.get("doi") or entry.get("DOI")
            document.chunks.append(
                Chunk(
                    source_id=document.id,
                    text=str(title),
                    kind="bibliography",
                    locator=Locator(json_path=f"entries[{index}]"),
                    citation_policy=CitationPolicy.VERIFY_FIRST,
                    metadata={"doi": doi} if doi else {},
                )
            )
        return ImportReport(source=document)


class EmailParser:
    name = "email"
    extensions = (".eml",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "message/rfc822")
        message = email.message_from_bytes(path.read_bytes(), policy=policy.default)
        body = message.get_body(preferencelist=("plain", "html"))
        content = body.get_content() if body else ""
        document.chunks.append(
            Chunk(
                source_id=document.id,
                text=str(content),
                kind="email",
                locator=Locator(message_id=message.get("Message-ID")),
                citation_policy=CitationPolicy.INTERNAL_ONLY,
                metadata={"subject": message.get("Subject", ""), "from": message.get("From", "")},
            )
        )
        return ImportReport(source=document)


class MboxParser:
    name = "mbox"
    extensions = (".mbox",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/mbox")
        seen: set[str] = set()
        for index, message in enumerate(mailbox.mbox(path)):
            message_id = message.get("Message-ID") or f"mbox:{index}"
            if message_id in seen:
                continue
            seen.add(message_id)
            payload = message.get_payload(decode=True)
            if isinstance(payload, bytes):
                text = payload.decode(message.get_content_charset() or "utf-8", errors="replace")
            else:
                text = str(payload or "")
            document.chunks.append(
                Chunk(
                    source_id=document.id,
                    text=text,
                    kind="email",
                    locator=Locator(message_id=message_id),
                    citation_policy=CitationPolicy.INTERNAL_ONLY,
                    metadata={
                        "subject": message.get("Subject", ""),
                        "date": message.get("Date", ""),
                    },
                )
            )
        return ImportReport(source=document)


class MsgParser:
    name = "msg"
    extensions = (".msg",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/vnd.ms-outlook")
        message = extract_msg.Message(path)  # type: ignore[no-untyped-call]
        try:
            document.chunks.append(
                Chunk(
                    source_id=document.id,
                    text=message.body or "",
                    kind="email",
                    locator=Locator(message_id=message.messageId),
                    citation_policy=CitationPolicy.INTERNAL_ONLY,
                    metadata={
                        "subject": message.subject or "",
                        "sender": message.sender or "",
                        "date": str(message.date or ""),
                    },
                )
            )
        finally:
            message.close()
        return ImportReport(source=document)


class EndNoteXmlParser:
    name = "endnote_xml"
    extensions = (".xml",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/xml")
        root = ElementTree.parse(path).getroot()
        records = root.findall(".//record")
        for index, record in enumerate(records):
            title = " ".join(text.strip() for text in record.itertext() if text.strip())
            doi_node = record.find(".//electronic-resource-num")
            doi = doi_node.text.strip() if doi_node is not None and doi_node.text else None
            if title:
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=title,
                        kind="bibliography",
                        locator=Locator(json_path=f"records[{index}]"),
                        citation_policy=CitationPolicy.VERIFY_FIRST,
                        metadata={"doi": doi} if doi else {},
                    )
                )
        if not records:
            raise ValueError("XML is not a supported EndNote export")
        return ImportReport(source=document)


class StructuredParser:
    name = "structured"
    extensions = (".json", ".ipynb")

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/json")
        if path.suffix == ".ipynb":
            notebook = nbformat.read(path, as_version=4)  # type: ignore[no-untyped-call]
            for index, cell in enumerate(notebook.cells):
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=cell.source,
                        kind=f"notebook_{cell.cell_type}",
                        locator=Locator(json_path=f"cells[{index}]"),
                        citation_policy=CitationPolicy.INTERNAL_ONLY,
                    )
                )
        else:
            data = json.loads(path.read_text(encoding="utf-8"))
            if isinstance(data, list) and all(isinstance(item, dict) for item in data):
                for index, item in enumerate(data):
                    text = str(item.get("content") or item.get("text") or item.get("message") or "")
                    if text:
                        document.chunks.append(
                            Chunk(
                                source_id=document.id,
                                text=text,
                                kind="chat_message",
                                locator=Locator(
                                    message_id=str(item.get("id") or index),
                                    json_path=f"$[{index}]",
                                ),
                                citation_policy=CitationPolicy.PROCESS_ONLY,
                                metadata={
                                    "participant": item.get("participant") or item.get("author"),
                                    "timestamp": item.get("timestamp") or item.get("time"),
                                },
                            )
                        )
            else:
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=json.dumps(data, ensure_ascii=False),
                        locator=Locator(json_path="$"),
                        citation_policy=CitationPolicy.INTERNAL_ONLY,
                    )
                )
        return ImportReport(source=document)


class CodeParser:
    name = "code"
    extensions = (".py", ".js", ".ts", ".tsx", ".java", ".c", ".cpp", ".rs", ".go", ".yaml", ".yml")

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "text/plain")
        text = path.read_text(encoding="utf-8", errors="replace")
        for number, line in enumerate(text.splitlines(), start=1):
            if line.strip():
                kind = (
                    "symbol" if re.match(r"\s*(def|class|function|interface)\s+", line) else "code"
                )
                secret_risk = bool(
                    re.search(r"(?i)(api[_-]?key|password|secret|token)\s*[:=]", line)
                )
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=line,
                        kind=kind,
                        locator=Locator(line_start=number, line_end=number),
                        citation_policy=CitationPolicy.INTERNAL_ONLY,
                        metadata={"credential_risk": secret_risk},
                    )
                )
        return ImportReport(source=document)


class ZipTreeParser:
    name = "archive"
    extensions = (".zip",)

    def parse(self, path: Path, digest: str) -> ImportReport:
        document = source(path, digest, self.name, "application/zip")
        with zipfile.ZipFile(path) as archive:
            if len(archive.infolist()) > 10_000:
                raise ValueError("Archive exceeds entry limit")
            for index, info in enumerate(archive.infolist()):
                member = Path(info.filename)
                if member.is_absolute() or ".." in member.parts:
                    raise ValueError("Archive contains path traversal")
                if info.file_size > 1_000_000_000:
                    raise ValueError("Archive member exceeds size limit")
                document.chunks.append(
                    Chunk(
                        source_id=document.id,
                        text=info.filename,
                        kind="archive_entry",
                        locator=Locator(json_path=f"entries[{index}]"),
                        citation_policy=CitationPolicy.INTERNAL_ONLY,
                        metadata={"size": info.file_size},
                    )
                )
        return ImportReport(source=document)


class ImageParser:
    name = "image"
    extensions = (".png", ".jpg", ".jpeg", ".webp", ".svg")

    def parse(self, path: Path, digest: str) -> ImportReport:
        media_type = (
            "image/svg+xml" if path.suffix == ".svg" else f"image/{path.suffix.lstrip('.')}"
        )
        document = source(path, digest, self.name, media_type)
        metadata: dict[str, object] = {}
        if path.suffix == ".svg":
            raw = path.read_text(encoding="utf-8", errors="replace")
            if re.search(r"(?i)<script|onload\s*=|javascript:", raw):
                raise ValueError("SVG contains active content")
            metadata["vector"] = True
        else:
            with Image.open(path) as image:
                metadata.update({"width": image.width, "height": image.height, "mode": image.mode})
        document.metadata.update(metadata)
        document.chunks.append(
            Chunk(
                source_id=document.id,
                text=path.name,
                kind="image_metadata",
                locator=Locator(json_path="$metadata"),
                citation_policy=CitationPolicy.INTERNAL_ONLY,
                metadata=metadata,
            )
        )
        return ImportReport(source=document)


def default_registry() -> IngestionRegistry:
    registry = IngestionRegistry()
    for parser in (
        TextParser(),
        HtmlParser(),
        PdfParser(),
        DocxParser(),
        PptxParser(),
        CsvParser(),
        XlsxParser(),
        BibliographyParser(),
        EmailParser(),
        MboxParser(),
        MsgParser(),
        EndNoteXmlParser(),
        StructuredParser(),
        CodeParser(),
        ZipTreeParser(),
        ImageParser(),
    ):
        registry.register(cast(Parser, parser))
    return registry
